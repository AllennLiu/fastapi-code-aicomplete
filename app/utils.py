import re
import torch
import string
import opencc
import pymupdf
import colorama
import markdown
import textwrap
import itertools
import asyncio
import aiofiles
import aiopathlib
import markdown.blockprocessors
from PIL import Image
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.worksheet import worksheet
from fastapi import UploadFile
from dataclasses import dataclass, field
from pygments.formatters import HtmlFormatter
from xml.etree.ElementTree import Element, SubElement
from typing import Any, Dict, List, Iterable, cast
from transformers import AutoTokenizer, AutoModelForCausalLM, PreTrainedTokenizer, PreTrainedTokenizerFast

tw_to_cn, cn_to_tw = opencc.OpenCC('tw2sp'), opencc.OpenCC('s2tw')

try:
    import chatglm_cpp
    enable_chatglm_cpp = True
except:
    print(textwrap.dedent("""\
    [WARN] chatglm-cpp not found. Install it by `pip install chatglm-cpp` for better performance.
    Check out https://github.com/li-plus/chatglm.cpp for more details.
    """))
    enable_chatglm_cpp = False

@dataclass
class ML:
    tokenizer : PreTrainedTokenizer | PreTrainedTokenizerFast | None = field()
    model     : chatglm_cpp.Pipeline = field()

def print_process(message: str) -> None:
    """
    Print colorful message with self-defined process arrow, this
    could make the message more readable.

    Args:
        message (str): Colorful message to print.
    """
    colorama.init(autoreset=True)
    print(f'{colorama.Fore.LIGHTCYAN_EX} ➠{colorama.Fore.RESET} {message}')

def load_llm(
    name             : str,
    quantize         : int,
    dtype            : str  = 'f32',
    transformers_only: bool = False,
    device           : str  = 'cpu',
    init_token       : bool = True,
    **_              : Any
) -> tuple[PreTrainedTokenizer | PreTrainedTokenizerFast | None, chatglm_cpp.Pipeline]:
    """
    使用 :module:`~chatglm_cpp` 量化加速推理，實現通過 ``CPU+MEM`` 也能
    與模型實時交互，取代原先通過 :module:`~transformers` 進行推理的方式

    - 量化精度：``fp32 > fp16 > int8 > int4`` 越小模型推理能力越差，但\
        對於硬體要求就越低，使用 ``CPU+MEM`` 建議設 ``fp16`` 或 ``int8``
    - 同上，``fp32`` 和 ``fp16`` 需在有 GPU 的設備上運行才能擁有正常\
        的推理表現

    Args:
        name (str): Model name.
        quantize (int): Quantization level.
        dtype (str): Data type. Defaults to 'f32'.
        transformers_only (bool): Whether to use transformers only. \
            Defaults to False.
        device (str): Load model physical device. Defaults to 'cpu'.
        init_token (bool): Whether to initialize tokenizer. \
            Defaults to True.

    Returns:
        tuple[PreTrainedTokenizer | PreTrainedTokenizerFast | None, chatglm_cpp.Pipeline]: \
            Tokenizer and pipeline instance.
    """
    print(f'Loading LLM {name} quantize: {quantize} dtype: {dtype} ({device}) ...')
    pretrained_args = dict(trust_remote_code=True, local_files_only=True, device=device)
    tokenizer = AutoTokenizer.from_pretrained(name, **pretrained_args) if init_token else None
    if not transformers_only:
        if enable_chatglm_cpp:
            print('Using chatglm-cpp to improve performance')
            if quantize in [ 4, 5, 8 ]:
                dtype = f'q{quantize}_0'
            pipeline = chatglm_cpp.Pipeline(name, dtype=dtype)
            return (tokenizer, pipeline)
        print('chatglm-cpp not enabled, falling back to transformers')
    pretrained_args = dict(trust_remote_code=True, local_files_only=True, low_cpu_mem_usage=True)
    model = AutoModelForCausalLM.from_pretrained(
        name, **pretrained_args, torch_dtype=torch.bfloat16)
    return (tokenizer, model.to(device).eval())

def copilot_prompt(lang_tags: Dict[str, Dict[str, str]], lang: str, prompt: str) -> tuple[str, str]:
    """
    依指定代碼語言來擷取注釋符號，連接 prompt 生成代碼描述，最後返回
    ``prompt 生成代碼描述`` 與 ``注釋符號`` 的元組

    >>> ('# language: Python\\n# 幫我寫一個冒泡排序\\n', '#')

    Args:
        lang_tags (Dict[str, Dict[str, str]]): 語言標籤映射字典
        lang (str): 代碼語言
        prompt (str): 代碼描述

    Returns:
        tuple[str, str]: 返回 ``prompt 生成代碼描述`` 與 ``注釋符號`` 的元組
    """
    pattern = re.compile('\s+language:.*')
    comment_char = pattern.sub('', lang_tags[lang]["mark"])
    result = f'{lang_tags[lang]["mark"]}\n{comment_char} {tw_to_cn.convert(prompt)}\n'
    return result, comment_char

def block_bad_words(content: str) -> str:
    """屏蔽在聊天信息中不允許的字符，全數替換成 ``Black.Milan``

    Args:
        content (str): 模型返回的文本信息

    Returns:
        str: 屏蔽關鍵字後的文本信息
    """
    pattern = re.compile('chatg[\w|-]+|清华大学.*KEG|智谱.*AI|GLM.*\d+', flags=re.I)
    return pattern.sub('Black.Milan', content)

def remove_punctuation(content: str) -> str:
    """
    將下列拼接起來使用 :func:`~str.maketrans` 建立翻譯表，將所有的
    標點符號映射成 `None`
    - 英文字符的符號 `string.punctuation`
    - 常用中文字符和符號的 `Unicode ASCII 碼` 範圍 ``0x3000`` 到 \
        ``0x303F``
    - 全角字符的 `Unicode 碼` 範圍 ``0xFF00`` 到 ``0xFFEF``

    最後使用 :func:`~str.translate` 刪除文本中所有的**中英文標點符號**

    Args:
        content (str): 任意文本

    Returns:
        str: 標點符號刪除後的文本

    Examples:
    >>> remove_punctuation('你好，世界！')
    '你好世界'
    """
    punctuation = string.punctuation
    punctuation += ''.join(chr(i) for i in itertools.chain(range(0x3000, 0x303F), range(0xFF00, 0xFFEF)))
    translator = str.maketrans('', '', punctuation)
    return re.sub(r'\s', '', content.translate(translator))

def set_similarity(array1: Iterable[str], array2: Iterable[str]) -> tuple[set, float]:
    """不同長度數組，計算 ``交集長度 / 並集長度 = 單詞數組相似度 (%)``

    Args:
        array1 (Iterable[str]): first iterable object
        array2 (Iterable[str]): second iterable object

    Returns:
        tuple[set, float]: 返回 `元組(交集, 相似度比例)`

    Examples:
    >>> set_similarity(['升级','测试', 'REDFISH'], ['用户', '命令', '升级', 'redfish'])
    ({'redfish', '升级'}, 40.0)
    """
    x, y = set(map(str.lower, array1)), set(map(str.lower, array2))
    intersections: set[str] = x.intersection(y)
    return intersections, round((len(intersections) / len(x.union(y))) * 100, 2)

async def read_file(file: UploadFile, uuid: str) -> str:
    """
    異步存取文件，暫存 ``.pdf``、``.docx``、``.pptx``、``.xlsx``
    文件為緩存至 `/tmp` 下，並依照對應文件副檔名讀取方法，轉換為
    文本後刪除

    Args:
        file (UploadFile): 上傳的 :class:`~fastapi.UploadFile` 對象
        uuid (str): 文件的唯一標識符

    Returns:
        str: 文件經解析後讀取到的文本內容
    """
    path = aiopathlib.AsyncPath(f'/tmp/{uuid}{file.filename.split(".")[-1]}')
    abs_path = str(path.absolute())
    await path.parent.mkdir(parents=True, exist_ok=True)
    async with aiofiles.open(path, 'wb') as out_file:
        buffer = cast(bytes, await file.read())
        await out_file.write(buffer)
    loop = asyncio.get_event_loop()
    if path.name.endswith('.pdf'):
        docs = await loop.run_in_executor(None, lambda: pymupdf.open(abs_path))
        content = ''.join(page.get_text() for page in docs).strip()
    elif path.name.endswith('.docx'):
        doc = await loop.run_in_executor(None, lambda: Document(abs_path))
        content = '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    elif path.name.endswith('.pptx'):
        content = ''
        present = await loop.run_in_executor(None, lambda: Presentation(abs_path))
        for slide in present.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    content += f'{shape.text}\n'
    elif path.name.endswith('.xlsx'):
        wb = await loop.run_in_executor(None, lambda: load_workbook(abs_path))
        worksheets = cast(worksheet.Worksheet, wb.active).values
        content = '\n'.join(','.join(map(str, filter(bool, row))) for row in worksheets)
    else:
        content = await path.read_text(encoding='utf-8', errors='ignore')
    await path.remove(missing_ok=True)
    return content.strip()

def verify_image(filename: str) -> bool:
    """檢查指定路徑文件是否為 :class:`~PIL.Image` 可識別的圖像格式

    Args:
        filename (str): image file

    Returns:
        bool: target is image or not
    """
    try:
        with Image.open(filename) as img:
            img.verify()
            return True
    except (IOError, SyntaxError):
        return False

class IgnoreOrderedListProcessor(markdown.blockprocessors.BlockProcessor):
    """
    Completely ignore serial table processors.
    """
    def test(self, parent: Element, block: str) -> bool:
        pattern_ordered_list = re.compile(r'^\s*\d+\.\s')
        return bool(pattern_ordered_list.match(block))

    def run(self, parent: Element, blocks: List[str]) -> None:
        block = blocks.pop(0)
        pre = SubElement(parent, 'remove')
        pre.text = block

class IgnoreOrderedListExtension(markdown.extensions.Extension):
    """
    Custom extension for handling Markdown files with complete
    ignore of serial tables.
    """
    def extendMarkdown(self, md: markdown.Markdown) -> None:
        if 'olist' in md.parser.blockprocessors:
            md.parser.blockprocessors.deregister('olist')
        md.parser.blockprocessors.register(
            item=IgnoreOrderedListProcessor(md.parser), name='ignore_ordered_list', priority=25.)

def md_to_html(content: str, html_formatter: bool = False) -> str:
    """轉換 ``Markdown`` 文本為 ``GitHub`` 編程風格的 HTML 文本

    Args:
        content (str): markdown 文本
        html_formatter (bool, optional): 需要轉換為 HTML 文本. \
            Defaults to False.

    Returns:
        str: 轉換後的 HTML 文本
    """
    md = markdown.Markdown(extensions=[ 'fenced_code', 'codehilite', IgnoreOrderedListExtension() ])
    html_content = md.convert(content)
    style = HtmlFormatter(style='github-dark').get_style_defs('.codehilite')
    pattern_remove = re.compile('</?remove>')
    return textwrap.dedent(f"""\
    {f'<head><style> {style} </style></head>' if html_formatter else ''}
    {pattern_remove.sub('', html_content)}
    """)

def md_no_codeblock(md: str) -> str:
    """Remove all codeblock markdown text and symbols.

    Args:
        md (str): markdown content.

    Returns:
        str: codeblock removed text.
    """
    pattern_codeblock = re.compile(r'`{3}\w{0,}(.*)`{3}', re.DOTALL)
    return pattern_codeblock.sub(r'\1', md)
